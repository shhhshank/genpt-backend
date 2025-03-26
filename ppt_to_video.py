import os
import re
import subprocess
import shutil
from typing import List, Dict
import google.generativeai as genai
from pptx import Presentation
from gtts import gTTS
from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance
import textwrap
import time
import numpy as np
import requests


class PPTVideoGenerator:
    def __init__(self, output_dir: str = "uploads/generated_videos"):
        self.output_dir = output_dir
        self.temp_dir = os.path.join("uploads", "temp_video")
        self.audio_dir = os.path.join(self.temp_dir, "audio")
        self.image_dir = os.path.join(self.temp_dir, "images")
        self.setup_directories()
        
        # Initialize Gemini API
    # Note: You'll need to set your API key in your environment or config
    API_KEY = "AIzaSyB1U3NavO8CvkQ2pk0NFpEf_NKYJPCnAPk"  # Replace with actual API key or load from environment
    print("Configuring Gemini API")
    genai.configure(api_key=API_KEY)
    
    # Configure the model
    generation_config = {
        "temperature": 0.4,
        "top_p": 1,
        "top_k": 32,
        "max_output_tokens": 2048,
    }
    
    # Create a Gemini model instance
    print("Creating Gemini model instance with gemini-1.5-pro")
    model = genai.GenerativeModel(
        model_name="gemini-1.5-pro",
        generation_config=generation_config
    )

    def setup_directories(self):
        """Ensure all required directories exist"""
        for dir_path in [self.output_dir, self.temp_dir, self.audio_dir, self.image_dir]:
            os.makedirs(dir_path, exist_ok=True)

    def extract_content_from_ppt(self, ppt_path: str) -> List[Dict]:
        prs = Presentation(ppt_path)
        slides_content = []
        
        for idx, slide in enumerate(prs.slides):
            content = {
                'index': idx,
                'title': '',
                'points': [],
                'image_path': None
            }
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text = shape.text.strip()
                    if not content['title']:
                        content['title'] = text
                    else:
                        content['points'].append(text)
                elif shape.shape_type == 13:  # Picture
                    image_path = os.path.join(self.image_dir, f"slide_{idx}.jpg")
                    os.makedirs(os.path.dirname(image_path), exist_ok=True)
                    with open(image_path, 'wb') as f:
                        f.write(shape.image.blob)
                    content['image_path'] = image_path
            
            # If no image was found, create an animated text slide (multiple frames)
            if not content['image_path']:
                print(f"Creating animated text slide for slide {idx + 1}")
                # Create base directory for animation frames
                frame_dir = os.path.join(self.image_dir, f"slide_{idx}_frames")
                os.makedirs(frame_dir, exist_ok=True)
                
                # Create animation frames
                self._create_animated_text_slide(
                    title=content['title'],
                    points=content['points'],
                    output_dir=frame_dir,
                    num_frames=5  # Adjust based on desired animation smoothness
                )
                
                # Set the first frame as the image path for concat file
                content['image_path'] = os.path.join(frame_dir, "frame_0.jpg")
                content['is_animated'] = True
                content['frame_dir'] = frame_dir
                content['frame_count'] = 5  # Same as num_frames above
            
            slides_content.append(content)
        
        return slides_content

    def _create_animated_text_slide(self, title: str, points: List[str], output_dir: str, num_frames: int = 30):
        """
        Create professional-grade animated text slides with modern design elements.
        Using 30fps for smooth animations.
        """
        try:
            from PIL import Image, ImageDraw, ImageFont, ImageFilter, ImageEnhance
            import numpy as np
            import math
            
            # Professional HD resolution
            width, height = 1920, 1080
            
            # Professional color scheme
            COLORS = {
                'primary': (41, 128, 185),      # Professional blue
                'secondary': (44, 62, 80),      # Dark slate
                'accent': (46, 204, 113),       # Emerald green
                'background': (236, 240, 241),  # Light gray
                'text': (52, 73, 94),           # Dark gray-blue
                'highlight': (52, 152, 219)     # Light blue
            }
            
            # Load professional fonts (with fallbacks)
            try:
                title_font = ImageFont.truetype("Montserrat-Bold.ttf", 120)
                subtitle_font = ImageFont.truetype("Montserrat-SemiBold.ttf", 80)
                body_font = ImageFont.truetype("Montserrat-Regular.ttf", 70)
            except:
                try:
                    title_font = ImageFont.truetype("Arial Bold.ttf", 120)
                    subtitle_font = ImageFont.truetype("Arial.ttf", 80)
                    body_font = ImageFont.truetype("Arial.ttf", 70)
                except:
                    print("⚠️ Falling back to default fonts")
                    title_font = ImageFont.load_default()
                    subtitle_font = title_font
                    body_font = title_font

            # Create frames for smooth animation
            for frame in range(num_frames):
                progress = frame / (num_frames - 1)
                
                # Create base image with gradient background
                background = Image.new('RGB', (width, height), COLORS['background'])
                
                # Add subtle pattern overlay
                pattern = self._create_pattern_overlay(width, height, progress)
                background.paste(pattern, (0, 0), pattern)
                
                # Add dynamic elements
                draw = ImageDraw.Draw(background)
                
                # Animated header bar
                header_height = int(200 + 20 * math.sin(progress * 2 * math.pi))
                self._draw_gradient_rectangle(
                    draw, 
                    [0, 0, width, header_height],
                    COLORS['primary'],
                    self._adjust_color(COLORS['primary'], 30)
                )
                
                # Title animation
                title_opacity = min(1.0, progress * 2)
                title_y = 300 + int(20 * math.sin(progress * 2 * math.pi))
                
                # Draw title with modern styling
                title_wrapped = self._wrap_text(title, title_font, width - 200)
                self._draw_text_with_effects(
                    draw,
                    background,
                    title_wrapped,
                    (width//2, title_y),
                    title_font,
                    progress,
                    is_title=True
                )
                
                # Animated separator
                separator_width = int(width * min(1.0, progress * 3))
                draw.line(
                    [(width//2 - separator_width//2, title_y + 100),
                     (width//2 + separator_width//2, title_y + 100)],
                    fill=COLORS['accent'],
                    width=5
                )
                
                # Animate bullet points
                y_position = title_y + 200
                for i, point in enumerate(points[:4]):  # Limit to 4 key points
                    point_progress = max(0, min(1.0, (progress * 3 - i * 0.2)))
                    if point_progress <= 0:
                        continue
                    
                    # Calculate dynamic positions
                    x_offset = int((1 - point_progress) * 200)
                    y_offset = int(10 * math.sin((progress + i * 0.2) * 2 * math.pi))
                    
                    # Draw modern bullet point
                    self._draw_bullet_point(
                        draw,
                        background,
                        point,
                        (width//6 + x_offset, y_position + y_offset),
                        body_font,
                        point_progress,
                        i
                    )
                    
                    y_position += 150
                
                # Add professional footer
                self._draw_footer(draw, background, height, progress)
                
                # Apply final effects
                background = self._apply_final_effects(background, progress)
                
                # Save frame with high quality
                frame_path = os.path.join(output_dir, f"frame_{frame:03d}.jpg")
                background.save(frame_path, 'JPEG', quality=95)
                
            return True
            
        except Exception as e:
            print(f"✗ Error creating animated slide: {str(e)}")
            return False

    def _create_pattern_overlay(self, width: int, height: int, progress: float) -> Image:
        """Create subtle, animated background pattern"""
        pattern = Image.new('RGBA', (width, height), (255, 255, 255, 0))
        draw = ImageDraw.Draw(pattern)
        
        # Create subtle grid pattern
        spacing = 50
        offset = int(progress * spacing)
        
        for x in range(-offset, width + spacing, spacing):
            for y in range(-offset, height + spacing, spacing):
                draw.ellipse([x-2, y-2, x+2, y+2], fill=(0, 0, 0, 10))
        
        return pattern

    def _draw_gradient_rectangle(self, draw, coords, color1, color2):
        """Draw a gradient-filled rectangle"""
        x1, y1, x2, y2 = coords
        for y in range(y1, y2):
            progress = (y - y1) / (y2 - y1)
            color = self._interpolate_color(color1, color2, progress)
            draw.line([(x1, y), (x2, y)], fill=color)

    def _draw_text_with_effects(self, draw, image, text, position, font, progress, is_title=False):
        """Draw text with professional effects"""
        x, y = position
        
        # Calculate text size
        text_width, text_height = draw.textsize(text, font=font)
        
        # Create mask for text effects
        mask = Image.new('L', (text_width + 40, text_height + 40), 0)
        mask_draw = ImageDraw.Draw(mask)
        
        # Draw text on mask
        mask_draw.text((20, 20), text, font=font, fill=255)
        
        # Apply blur for glow effect
        glow = mask.filter(ImageFilter.GaussianBlur(10))
        
        # Composite glow onto main image
        image.paste(
            (255, 255, 255),
            (int(x - text_width/2 - 20), int(y - text_height/2 - 20)),
            glow
        )
        
        # Draw main text
        draw.text(
            (x, y),
            text,
            font=font,
            fill=(52, 73, 94) if is_title else (44, 62, 80),
            anchor="mm"
        )

    def _draw_bullet_point(self, draw, image, text, position, font, progress, index):
        """Draw an animated bullet point with modern styling"""
        x, y = position
        
        # Bullet point styling
        bullet_size = int(20 * progress)
        bullet_color = self._interpolate_color(
            (41, 128, 185),
            (46, 204, 113),
            index / 4
        )
        
        # Draw bullet
        draw.ellipse(
            [x - bullet_size, y - bullet_size, x + bullet_size, y + bullet_size],
            fill=bullet_color
        )
        
        # Draw text with shadow
        text_x = x + bullet_size + 40
        shadow_offset = 3
        
        # Shadow
        draw.text(
            (text_x + shadow_offset, y + shadow_offset),
            text,
            font=font,
            fill=(0, 0, 0, 50),
            anchor="lm"
        )
        
        # Main text
        draw.text(
            (text_x, y),
            text,
            font=font,
            fill=(52, 73, 94),
            anchor="lm"
        )

    def _draw_footer(self, draw, image, height, progress):
        """Draw professional footer with animation"""
        footer_height = 80
        y = height - footer_height
        
        # Draw footer background
        self._draw_gradient_rectangle(
            draw,
            [0, y, image.width, height],
            (41, 128, 185, int(255 * progress)),
            (52, 152, 219, int(255 * progress))
        )

    def _apply_final_effects(self, image, progress):
        """Apply final professional effects to the frame"""
        # Enhance contrast slightly
        enhancer = ImageEnhance.Contrast(image)
        image = enhancer.enhance(1.1)
        
        # Add subtle vignette
        vignette = self._create_vignette(image.size, progress)
        image.paste(vignette, (0, 0), vignette)
        
        return image

    def _create_vignette(self, size, progress):
        """Create subtle vignette effect"""
        vignette = Image.new('RGBA', size, (0, 0, 0, 0))
        draw = ImageDraw.Draw(vignette)
        
        width, height = size
        for i in range(50):
            opacity = int(i * progress * 3)
            box = [i, i, width - i, height - i]
            draw.rectangle(box, outline=(0, 0, 0, opacity))
        
        return vignette

    def _interpolate_color(self, color1, color2, progress):
        """Smoothly interpolate between two colors"""
        return tuple(
            int(c1 + (c2 - c1) * progress)
            for c1, c2 in zip(color1, color2)
        )

    def _adjust_color(self, color, adjustment):
        """Adjust color brightness"""
        return tuple(
            min(255, max(0, c + adjustment))
            for c in color
        )

    def _wrap_text(self, text: str, font, max_width: int) -> str:
        """Intelligently wrap text to fit width"""
        words = text.split()
        lines = []
        current_line = []
        
        for word in words:
            current_line.append(word)
            line = ' '.join(current_line)
            w, _ = font.getsize(line)
            
            if w > max_width:
                if len(current_line) == 1:
                    lines.append(current_line.pop())
                else:
                    current_line.pop()
                    lines.append(' '.join(current_line))
                    current_line = [word]
        
        if current_line:
            lines.append(' '.join(current_line))
        
        return '\n'.join(lines)

    def generate_script(self, slides_content: List[Dict]) -> List[Dict]:
        """Generate more engaging scripts using Gemini"""
        scripts = []
        
        for slide in slides_content:
            prompt = f"""
            Create an engaging, professional narration script for a presentation slide.
            Make it sound like a skilled presenter or TED talk speaker.
            
            Title: {slide['title']}
            Points: {slide['points']}
            
            Requirements:
            1. Use natural, conversational language
            2. Include brief pauses (marked with '...')
            3. Keep it concise (20-30 seconds when spoken)
            4. Make it engaging and dynamic
            5. Add emphasis on key words (in *asterisks*)
            6. Use transitions between points
            7. End with a strong concluding statement
            
            Format the output as:
            <<script>>
            
            Example:
            <<Let's explore the fascinating world of *artificial intelligence*... 
            At its core, AI is revolutionizing how we approach problem-solving... 
            What's particularly exciting is how it's transforming healthcare, 
            making diagnoses more accurate and treatments more effective...>>
            """
            
            response = self.model.generate_content(prompt)
            script = self._extract_prompt(response.text)
            
            scripts.append({
                'slide_index': slide['index'],
                'script': script
            })
            
        return scripts

    def _extract_prompt(self, text: str) -> str:
        """Extract text between << and >> markers"""
        pattern = r"<<(.+?)>>"
        match = re.search(pattern, text, re.DOTALL)
        return match.group(1) if match else text

    def generate_audio(self, scripts: List[Dict]) -> List[str]:
        """Generate audio using ElevenLabs API with gTTS fallback"""
        try:
            audio_paths = []
            ELEVENLABS_VOICE_ID = "JBFqnCBsd6RMkjVDRZzb"  # Antoni voice ID
            API_KEY = "sk_4fb2de074dd4b122bf1e3d5c510f4ea7f33ab5c6fcc600f0"
            
            for script in scripts:
                try:
                    audio_path = os.path.join(self.audio_dir, f"slide_{script['slide_index']}.mp3")
                    
                    # ElevenLabs API endpoint
                    url = f"https://api.elevenlabs.io/v1/text-to-speech/{ELEVENLABS_VOICE_ID}"
                    
                    headers = {
                        "xi-api-key": API_KEY,
                        "Content-Type": "application/json"
                    }
                    
                    data = {
                        "text": script['script'],
                        "model_id": "eleven_multilingual_v2"
                    }
                    
                    # Make API request
                    response = requests.post(
                        f"{url}?output_format=mp3_44100_128",
                        headers=headers,
                        json=data
                    )
                    
                    if response.status_code == 200:
                        # Save the audio file
                        with open(audio_path, 'wb') as f:
                            f.write(response.content)
                        print(f"✓ Generated audio with ElevenLabs for slide {script['slide_index'] + 1}")
                    else:
                        print(f"⚠️ ElevenLabs API error: {response.text}")
                        print("Falling back to gTTS for this slide...")
                        audio_path = self._generate_single_audio_gtts(script)
                    
                    audio_paths.append(audio_path)
                    time.sleep(0.5)  # Rate limiting
                    
                except Exception as e:
                    print(f"⚠️ Error generating audio for slide {script['slide_index'] + 1}: {str(e)}")
                    print("Falling back to gTTS for this slide...")
                    audio_path = self._generate_single_audio_gtts(script)
                    audio_paths.append(audio_path)
            
            return audio_paths
            
        except Exception as e:
            print(f"✗ Error in audio generation: {str(e)}")
            return self._generate_audio_fallback(scripts)

    def _generate_single_audio_gtts(self, script: Dict) -> str:
        """Generate a single audio file using gTTS"""
        try:
            audio_path = os.path.join(self.audio_dir, f"slide_{script['slide_index']}.mp3")
            
            # Add pauses for better pacing
            text_with_pauses = self._add_speech_pauses(script['script'])
            
            # Generate audio with gTTS
            tts = gTTS(text=text_with_pauses, lang='en', slow=False)
            tts.save(audio_path)
            
            print(f"✓ Generated audio with gTTS for slide {script['slide_index'] + 1}")
            return audio_path
            
        except Exception as e:
            print(f"✗ Error generating gTTS audio: {str(e)}")
            raise

    def _add_speech_pauses(self, text: str) -> str:
        """Add natural pauses to the text for better gTTS output"""
        # Add pause after sentences
        text = text.replace('. ', '... ')
        
        # Add pause after commas
        text = text.replace(', ', '... ')
        
        # Add pause for emphasis
        text = text.replace('*', '... ')
        
        return text

    def _generate_audio_fallback(self, scripts: List[Dict]) -> List[str]:
        """Generate all audio using gTTS"""
        print("⚠️ Using gTTS for all audio generation")
        audio_paths = []
        
        for script in scripts:
            try:
                audio_path = self._generate_single_audio_gtts(script)
                audio_paths.append(audio_path)
            except Exception as e:
                print(f"✗ Error in gTTS fallback for slide {script['slide_index'] + 1}: {str(e)}")
                # Create silent audio as last resort
                audio_path = self._generate_silent_audio(3.0)  # 3 seconds of silence
                audio_paths.append(audio_path)
        
        return audio_paths

    def _generate_silent_audio(self, duration: float) -> str:
        """Generate silent audio file as last resort"""
        try:
            audio_path = os.path.join(self.audio_dir, f"silence_{int(time.time())}.mp3")
            
            # Generate silent audio using FFmpeg
            subprocess.run([
                'ffmpeg', '-y',
                '-f', 'lavfi',
                '-i', f'anullsrc=r=44100:cl=stereo:d={duration}',
                '-c:a', 'libmp3lame',
                '-q:a', '2',
                audio_path
            ], check=True)
            
            return audio_path
        except Exception as e:
            print(f"✗ Error generating silent audio: {str(e)}")
            raise

    def create_video(self, slides_content: List[Dict], audio_paths: List[str]) -> str:
        """Create high-quality video with professional transitions"""
        try:
            # First convert all MP3s to WAV format for consistent audio handling
            wav_paths = []
            for idx, audio_path in enumerate(audio_paths):
                wav_path = os.path.join(self.temp_dir, f"audio_{idx}.wav")
                subprocess.run([
                    'ffmpeg', '-y',
                    '-i', audio_path,
                    '-acodec', 'pcm_s16le',
                    '-ar', '44100',
                    wav_path
                ], check=True)
                wav_paths.append(wav_path)

            # Combine WAV files
            combined_audio = os.path.join(self.temp_dir, "combined_audio.wav")
            audio_list = os.path.join(self.temp_dir, "audio_list.txt")
            
            with open(audio_list, 'w', encoding='utf-8') as f:
                for wav_path in wav_paths:
                    f.write(f"file '{os.path.abspath(wav_path)}'\n")
            
            # Combine audio files
            subprocess.run([
                'ffmpeg', '-y',
                '-f', 'concat',
                '-safe', '0',
                '-i', audio_list,
                '-c:a', 'pcm_s16le',
                combined_audio
            ], check=True)
            
            # Create video input file
            input_file = os.path.join(self.temp_dir, "input.txt")
            with open(input_file, 'w', encoding='utf-8') as f:
                for idx, (slide, audio) in enumerate(zip(slides_content, wav_paths)):
                    duration = self._get_audio_duration(audio)
                    
                    if slide.get('is_animated', False):
                        frame_dir = slide['frame_dir']
                        frame_duration = duration / 30  # 30fps
                        
                        for frame in range(30):
                            frame_path = os.path.join(frame_dir, f"frame_{frame:03d}.jpg")
                            if os.path.exists(frame_path):
                                f.write(f"file '{os.path.abspath(frame_path)}'\n")
                                f.write(f"duration {frame_duration}\n")
                    else:
                        image_path = os.path.abspath(slide['image_path'])
                        if os.path.exists(image_path):
                            f.write(f"file '{image_path}'\n")
                            f.write(f"duration {duration}\n")
            
            # Create final video with proper audio encoding
            output_path = os.path.join(
                self.output_dir,
                f"presentation_video_{int(time.time())}.mp4"
            )
            
            subprocess.run([
                'ffmpeg', '-y',
                '-f', 'concat',
                '-safe', '0',
                '-i', input_file,
                '-i', combined_audio,
                '-c:v', 'libx264',
                '-preset', 'slow',
                '-crf', '23',
                '-c:a', 'aac',
                '-b:a', '192k',
                '-pix_fmt', 'yuv420p',
                '-movflags', '+faststart',
                '-vf', 'fps=30',
                '-shortest',
                output_path
            ], check=True)
            
            print("✓ Created high-quality video presentation")
            return output_path
            
        except subprocess.CalledProcessError as e:
            print(f"✗ FFmpeg error: {e.output if hasattr(e, 'output') else str(e)}")
            raise
        except Exception as e:
            print(f"✗ Error creating video: {str(e)}")
            raise

    def _get_audio_duration(self, audio_path: str) -> float:
        """Get duration of audio file using FFprobe"""
        try:
            result = subprocess.run([
                'ffprobe',
                '-v', 'error',
                '-show_entries', 'format=duration',
                '-of', 'default=noprint_wrappers=1:nokey=1',
                audio_path
            ], capture_output=True, text=True, check=True)
            
            duration = float(result.stdout.strip())
            return duration
        except subprocess.CalledProcessError as e:
            print(f"✗ FFprobe error: {str(e)}")
            return 3.0  # Default duration if unable to get actual duration
        except Exception as e:
            print(f"✗ Error getting audio duration: {str(e)}")
            return 3.0

    def cleanup(self):
        """Clean up temporary files after processing"""
        try:
            if os.path.exists(self.temp_dir):
                shutil.rmtree(self.temp_dir)
                print("✓ Temporary files cleaned up")
        except Exception as e:
            print(f"✗ Error during cleanup: {str(e)}")