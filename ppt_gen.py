from pptx import Presentation
import io
import re
import inspect
import copy
import time
from PIL import Image
from pptx.util import Inches

from pptx.dml.color import _NoneColor

def sanitize_string(input_str):
    sanitized = re.sub(r"[^A-Za-z0-9_.-]", "", input_str)
    sanitized = re.sub(r"\.{2,}", ".", sanitized)
    sanitized = re.sub(r"^[^A-Za-z0-9]+", "", sanitized)
    sanitized = re.sub(r"[^A-Za-z0-9]+$", "", sanitized)
    sanitized = sanitized[:63] if len(sanitized) > 63 else sanitized.ljust(3, "_")
    return sanitized



def duplicate_slide(pres, template, idx):
    # Add a new slide
    copied_slide = pres.slides.add_slide(template.slide_layout)
    
    # Delete the existing shapes that are part of the layout
    for shp in copied_slide.shapes:
        copied_slide.shapes.element.remove(shp.element)        
    
    # Perform a deep copy of the shapes from the template
    for shp in template.shapes:
        if "Picture" in shp.name:
            img = io.BytesIO(shp.image.blob)
            copied_slide.shapes.add_picture(image_file = img,
                                            left = shp.left,
                                            top = shp.top,
                                            width = shp.width,
                                            height = shp.height)
        else:
            el = shp.element
            newel = copy.deepcopy(el)
            copied_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    slides = pres.slides._sldIdLst  
    slides.insert(idx, slides[-1])  # Move the last slide to the specified index

    return copied_slide
    


def duplicate_run_in_paragraph(paragraph, run_to_duplicate):

    # Create a new Run in the same Paragraph
    new_run = paragraph.add_run()


    # Copy font properties (if set)
    try:
        if run_to_duplicate.font.size:
            new_run.font.size = run_to_duplicate.font.size
        if run_to_duplicate.font.bold is not None:
            new_run.font.bold = run_to_duplicate.font.bold
        if run_to_duplicate.font.italic is not None:
            new_run.font.italic = run_to_duplicate.font.italic

        if run_to_duplicate.font.color and run_to_duplicate.font.color.rgb:
                new_run.font.color.rgb = run_to_duplicate.font.color.rgb
        if run_to_duplicate.font.name:
            new_run.font.name = run_to_duplicate.font.name
    except:
        pass
        


    return new_run

def duplicate_bullet(paragraph, shape):

    # Duplicate the paragraph within the same shape
    new_paragraph = shape.text_frame.add_paragraph()
    new_paragraph.text = paragraph.text
    new_paragraph.font.size = paragraph.font.size
    new_paragraph.font.bold = paragraph.font.bold
    new_paragraph.font.italic = paragraph.font.italic
    new_paragraph.font.underline = paragraph.font.underline
    new_paragraph.level = paragraph.level
    new_paragraph.bullet = True
    
    print(inspect.getmembers(paragraph))


    return new_paragraph

def replace_text(slide, search_text, replace_text):

    # Create a case-insensitive pattern for search_text
    pattern = re.compile(re.escape(search_text), re.IGNORECASE)

    for shape in slide.shapes:
        # Check if the shape has a text frame (contains text)
        if not shape.has_text_frame:
            continue
        
        # Loop through paragraphs and runs within the text frame
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                # Replace text using the case-insensitive pattern
                if re.search(pattern, run.text):
                    run.text = re.sub(pattern, replace_text, run.text)
    
    return slide

def get_view_from_text(slide, search_text):
    """Get paragraph containing exact search text"""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        
        for paragraph in shape.text_frame.paragraphs:
            # Remove debug print that was causing confusion
            if paragraph.text.strip() == search_text.strip():
                return paragraph
    
    print(f"✗ No shape found with text: {search_text}")
    return None


def fill_title(ppt, title, subtitle):
    slide = ppt.slides[0]
    replace_text(slide, "[title]", title)
    replace_text(slide, "[subtitle]", subtitle)

def fill_overview(ppt, overview):
    title_search = "[slide_title_i]"
    slide = ppt.slides[1]
    paragraph = get_view_from_text(slide, title_search)

    if paragraph:
        run_to_dup = paragraph.runs[0]
        paragraph.text = ""
        idx = -1
        for item in overview:
            idx += 1
            r_dup = duplicate_run_in_paragraph(paragraph, run_to_dup)
            r_dup.text = item
            if (idx < len(overview) - 1):
                r_dup.text += '\n'
    else:
        print("No paragraph found for: " + title_search)    

def process_image_for_shape(image_path, target_width, target_height):
    """Process image to match shape dimensions while maintaining aspect ratio"""
    try:
        # Convert EMU to pixels (914400 EMU = 1 inch, assuming 96 DPI)
        target_width_px = int(target_width / 914400 * 96)
        target_height_px = int(target_height / 914400 * 96)
        
        print(f"Processing image to {target_width_px}x{target_height_px} pixels")
        
        with Image.open(image_path) as img:
            # Convert to RGB if necessary
            if img.mode in ('RGBA', 'P'):
                img = img.convert('RGB')
            
            print(f"Original image size: {img.width}x{img.height}")
            
            # Calculate aspect ratios
            img_aspect = img.width / img.height
            target_aspect = target_width_px / target_height_px
            
            # Crop image to match target aspect ratio
            if img_aspect > target_aspect:
                # Image is wider than needed
                new_width = int(img.height * target_aspect)
                left = (img.width - new_width) // 2
                img = img.crop((left, 0, left + new_width, img.height))
            else:
                # Image is taller than needed
                new_height = int(img.width / target_aspect)
                top = (img.height - new_height) // 2
                img = img.crop((0, top, img.width, top + new_height))
            
            print(f"After cropping: {img.width}x{img.height}")
            
            # Resize to exact dimensions
            img = img.resize((target_width_px, target_height_px), Image.LANCZOS)
            
            print(f"After resizing: {img.width}x{img.height}")
            
            # Save to bytes
            img_byte_arr = io.BytesIO()
            img.save(img_byte_arr, format='JPEG', quality=85)
            img_byte_arr.seek(0)
            
            print("✓ Image processed successfully")
            return img_byte_arr
    except Exception as e:
        print(f"✗ Error processing image {image_path}: {str(e)}")
        return None

def fill_content(ppt, content):
    slide = ppt.slides[2]
    
    print(f"Duplicating {len(content) - 1} slides")
    for i in range(len(content) - 1):
        duplicate_slide(ppt, slide, i + 3)
    
    idx = 2
    for item in content:
        print(f"\nProcessing slide {idx + 1}")
        curr_slide = ppt.slides[idx]
        
        # Handle title
        title_para = get_view_from_text(curr_slide, "[slide_top_title_i]")
        if title_para:
            title_para.text = item['title']
            print(f"✓ Added title: {item['title']}")
        
        # Handle bullet points
        bullet_para = get_view_from_text(curr_slide, "[slide_i_point_j]")
        if bullet_para:
            run_to_dup = bullet_para.runs[0]
            bullet_para.text = ""
            pt_idx = -1
            for point in item['points']:
                pt_idx += 1
                r_dup = duplicate_run_in_paragraph(bullet_para, run_to_dup)
                r_dup.text = point
                if (pt_idx < len(item['points']) - 1):
                    r_dup.text += '\n'
            print(f"✓ Added {len(item['points'])} points")
        else:
            print("✗ No bullet point placeholder found")
        
        # Handle image if present
        if 'image' in item and item['image']:
            image_placeholder = None
            for shape in curr_slide.shapes:
                if shape.has_text_frame and "[image_i]" in shape.text:
                    image_placeholder = shape
                    break
            
            if image_placeholder:
                try:
                    # Store original properties
                    left = image_placeholder.left
                    top = image_placeholder.top
                    width = image_placeholder.width
                    height = image_placeholder.height
                    
                    # Get the desired z-order before removing the placeholder
                    desired_z_order = list(curr_slide.shapes).index(image_placeholder)
                    
                    # Process and resize image
                    processed_image = process_image_for_shape(
                        item['image'],
                        width,
                        height
                    )
                    
                    if processed_image:
                        # Remove the placeholder shape
                        sp = image_placeholder.element
                        sp.getparent().remove(sp)
                        
                        # Add the new image
                        picture = curr_slide.shapes.add_picture(
                            processed_image,
                            left,
                            top,
                            width,
                            height
                        )
                        
                        # Fix z-order by moving the picture element to the correct position
                        picture_element = picture.element
                        spTree = picture_element.getparent()
                        spTree.remove(picture_element)
                        
                        # Get all shape elements
                        shape_elements = spTree.findall('.//{*}sp') + spTree.findall('.//{*}pic')
                        
                        # Insert at the correct position
                        if desired_z_order < len(shape_elements):
                            ref_element = shape_elements[desired_z_order]
                            spTree.insert(list(spTree).index(ref_element), picture_element)
                        else:
                            # If it should be on top, append to end
                            spTree.append(picture_element)
                        
                        print(f"✓ Added image to slide {idx + 1} at z-index {desired_z_order}")
                    else:
                        print(f"✗ Failed to process image for slide {idx + 1}")
                except Exception as e:
                    print(f"✗ Error adding image to slide {idx + 1}: {str(e)}")
            else:
                print("✗ No image placeholder found")
        
        idx += 1
        print(f"Completed slide {idx}")

def ppt_gen(template_path, slide_data):
    ppt = Presentation(template_path)

    title = slide_data['title']
    subtitle = slide_data['subtitle']
    overview = slide_data['overview']
    content = slide_data['content']

    fill_title(ppt, title, subtitle)
    fill_overview(ppt, overview)
    fill_content(ppt, content)



    path = "uploads/generated_ppt/" + str(time.time() * 1000) +".pptx"
    ppt.save(path)


    return  '.' + path 
