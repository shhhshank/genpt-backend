import os
from typing import Dict, List
import google.generativeai as genai
from pptx import Presentation
from docx import Document
from docx.shared import Inches, Pt
from docx.enum.text import WD_ALIGN_PARAGRAPH
import time
from flask import Blueprint, request, jsonify, send_file
from werkzeug.utils import secure_filename
import subprocess



class MinutesGenerator:
    def __init__(self):
        # Initialize Gemini API
        API_KEY = "AIzaSyB1U3NavO8CvkQ2pk0NFpEf_NKYJPCnAPk"
        genai.configure(api_key=API_KEY)
        
        # Configure the model
        generation_config = {
            "temperature": 0.4,
            "top_p": 1,
            "top_k": 32,
            "max_output_tokens": 2048,
        }
        
        self.model = genai.GenerativeModel(
            model_name="gemini-1.5-pro",
            generation_config=generation_config
        )
        
        # Create output directory
        self.output_dir = "uploads/generated_minutes"
        os.makedirs(self.output_dir, exist_ok=True)

    def extract_ppt_content(self, ppt_path: str) -> List[Dict]:
        """Extract content from PowerPoint slides"""
        presentation = Presentation(ppt_path)
        slides_content = []
        
        for idx, slide in enumerate(presentation.slides):
            content = {
                'slide_number': idx + 1,
                'title': '',
                'points': [],
                'notes': ''
            }
            
            # Extract title
            if slide.shapes.title:
                content['title'] = slide.shapes.title.text
            
            # Extract points and text
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    if shape.text.strip() and shape.text != content['title']:
                        content['points'].append(shape.text.strip())
            
            # Extract notes if available
            if slide.notes_slide and slide.notes_slide.notes_text_frame:
                content['notes'] = slide.notes_slide.notes_text_frame.text
            
            slides_content.append(content)
        
        return slides_content

    def generate_minutes(self, slides_content: List[Dict], custom_prompt: str = None) -> str:
        """Generate meeting minutes using Gemini"""
        base_prompt = """
        Generate detailed meeting minutes from the following presentation content.
        Format the output as a structured document with:
        1. Meeting Overview
        2. Key Discussion Points
        3. Decisions Made
        4. Action Items
        5. Next Steps
        
        Use professional language and maintain a clear, organized structure.
        """
        
        if custom_prompt:
            base_prompt = f"{base_prompt}\nAdditional requirements: {custom_prompt}"
        
        content_text = "Presentation Content:\n"
        for slide in slides_content:
            content_text += f"\nSlide {slide['slide_number']}: {slide['title']}\n"
            content_text += "Points:\n" + "\n".join([f"- {point}" for point in slide['points']])
            if slide['notes']:
                content_text += f"\nNotes: {slide['notes']}\n"
        
        try:
            response = self.model.generate_content(base_prompt + "\n\n" + content_text)
            return response.text
        except Exception as e:
            print(f"Error generating minutes: {str(e)}")
            return None

    def create_document(self, minutes_content: str, format: str = 'docx') -> str:
        """Create a document with the generated minutes"""
        try:
            timestamp = int(time.time())
            output_path = None
            
            if format.lower() == 'docx':
                doc = Document()
                
                # Add title
                title = doc.add_heading('Meeting Minutes', 0)
                title.alignment = WD_ALIGN_PARAGRAPH.CENTER
                
                # Add date
                date_paragraph = doc.add_paragraph()
                date_paragraph.alignment = WD_ALIGN_PARAGRAPH.RIGHT
                date_paragraph.add_run(time.strftime("%B %d, %Y"))
                
                # Add content with formatting
                sections = minutes_content.split('\n\n')
                for section in sections:
                    if section.strip():
                        if section.startswith('#'):
                            # Handle headings
                            level = min(len(section.split()[0].strip('#')), 9)
                            text = ' '.join(section.split()[1:])
                            doc.add_heading(text, level)
                        else:
                            # Regular paragraph
                            para = doc.add_paragraph()
                            # Check for bullet points
                            if section.strip().startswith('•') or section.strip().startswith('-'):
                                para.style = 'List Bullet'
                            para.add_run(section.strip())
                
                output_path = os.path.join(self.output_dir, f"minutes_{timestamp}.docx")
                doc.save(output_path)
                print(f"✓ Created DOCX document: {output_path}")
                
            elif format.lower() == 'pdf':
                # First create DOCX
                docx_path = self.create_document(minutes_content, 'docx')
                output_path = os.path.join(self.output_dir, f"minutes_{timestamp}.pdf")
                
                # Try multiple PDF conversion methods
                converted = False
                
                # Method 1: Try docx2pdf
                if not converted:
                    try:
                        from docx2pdf import convert
                        convert(docx_path, output_path)
                        converted = True
                        print("✓ Converted to PDF using docx2pdf")
                    except ImportError:
                        print("⚠️ docx2pdf not available")
                    except Exception as e:
                        print(f"⚠️ docx2pdf conversion failed: {str(e)}")
                
                # Method 2: Try win32com (Windows only)
                if not converted:
                    try:
                        import win32com.client
                        word = win32com.client.Dispatch('Word.Application')
                        doc = word.Documents.Open(os.path.abspath(docx_path))
                        doc.SaveAs(os.path.abspath(output_path), FileFormat=17)  # 17 = PDF
                        doc.Close()
                        word.Quit()
                        converted = True
                        print("✓ Converted to PDF using Word COM object")
                    except ImportError:
                        print("⚠️ win32com not available")
                    except Exception as e:
                        print(f"⚠️ Word COM conversion failed: {str(e)}")
                
                # Method 3: Try LibreOffice (Linux/Mac)
                if not converted:
                    try:
                        subprocess.run([
                            'libreoffice',
                            '--headless',
                            '--convert-to',
                            'pdf',
                            '--outdir',
                            os.path.dirname(output_path),
                            docx_path
                        ], check=True)
                        converted = True
                        print("✓ Converted to PDF using LibreOffice")
                    except FileNotFoundError:
                        print("⚠️ LibreOffice not available")
                    except Exception as e:
                        print(f"⚠️ LibreOffice conversion failed: {str(e)}")
                
                # If all PDF conversion methods fail, return DOCX
                if not converted:
                    print("⚠️ All PDF conversion methods failed, returning DOCX instead")
                    output_path = docx_path
                else:
                    # Clean up DOCX if PDF conversion succeeded
                    try:
                        if os.path.exists(docx_path):
                            os.remove(docx_path)
                    except:
                        pass
            
            if not output_path or not os.path.exists(output_path):
                raise ValueError("Failed to create output document")
            
            return output_path
            
        except Exception as e:
            print(f"✗ Error creating document: {str(e)}")
            raise
